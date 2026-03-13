from __future__ import annotations

from dataclasses import dataclass
import json
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "docs" / "presentations"
OUTPUT_PATH = OUTPUT_DIR / "shelf-ai-plugin-architecture-cn.pptx"


@dataclass(frozen=True)
class PresentationAssets:
    package_json_path: Path
    hero_image: Path
    workbench_image: Path
    chat_image: Path
    branding_image: Path


@dataclass(frozen=True)
class PresentationTheme:
    slide_width: Any
    slide_height: Any
    bg: RGBColor
    paper: RGBColor
    ink: RGBColor
    muted: RGBColor
    accent: RGBColor
    accent_alt: RGBColor
    accent_dark: RGBColor
    gold: RGBColor
    soft: RGBColor
    ok: RGBColor
    warn: RGBColor
    code_bg: RGBColor
    title_font: str
    body_font: str
    mono_font: str


@dataclass(frozen=True)
class PluginStats:
    version: str
    command_count: int
    config_count: int


ASSETS = PresentationAssets(
    package_json_path=ROOT / "tools" / "vscode" / "shelf-ai" / "package.json",
    hero_image=ROOT / "docs" / "github-social-preview.png",
    workbench_image=ROOT / "docs" / "verification" / "knowledge-base-workbench.png",
    chat_image=ROOT / "docs" / "verification" / "knowledge-base-chat-first.png",
    branding_image=ROOT / "docs" / "branding" / "shelf-rainbow-lockup-options.png",
)

THEME = PresentationTheme(
    slide_width=Inches(13.333),
    slide_height=Inches(7.5),
    bg=RGBColor(247, 244, 238),
    paper=RGBColor(255, 255, 255),
    ink=RGBColor(34, 38, 46),
    muted=RGBColor(91, 96, 107),
    accent=RGBColor(239, 92, 58),
    accent_alt=RGBColor(27, 117, 119),
    accent_dark=RGBColor(42, 60, 66),
    gold=RGBColor(234, 178, 74),
    soft=RGBColor(240, 235, 226),
    ok=RGBColor(63, 133, 113),
    warn=RGBColor(212, 125, 34),
    code_bg=RGBColor(34, 38, 46),
    title_font="Aptos Display",
    body_font="Aptos",
    mono_font="Cascadia Code",
)

SLIDE_WIDTH = THEME.slide_width
SLIDE_HEIGHT = THEME.slide_height

BG = THEME.bg
PAPER = THEME.paper
INK = THEME.ink
MUTED = THEME.muted
ACCENT = THEME.accent
ACCENT_ALT = THEME.accent_alt
ACCENT_DARK = THEME.accent_dark
GOLD = THEME.gold
SOFT = THEME.soft
OK = THEME.ok
WARN = THEME.warn
CODE_BG = THEME.code_bg

TITLE_FONT = THEME.title_font
BODY_FONT = THEME.body_font
MONO_FONT = THEME.mono_font


def rgb_to_str(color: RGBColor) -> str:
    return f"{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def load_package_json() -> dict[str, Any]:
    return json.loads(ASSETS.package_json_path.read_text(encoding="utf-8"))


def build_plugin_stats(package_json: dict[str, Any]) -> PluginStats:
    return PluginStats(
        version=str(package_json.get("version", "unknown")),
        command_count=len(package_json.get("contributes", {}).get("commands", [])),
        config_count=len(package_json.get("contributes", {}).get("configuration", {}).get("properties", {})),
    )


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, radius: bool = True):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_outline_rect(
    slide,
    left,
    top,
    width,
    height,
    line_color: RGBColor,
    fill_color: RGBColor = PAPER,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.2)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int,
    color: RGBColor = INK,
    bold: bool = False,
    font_name: str = BODY_FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines: Iterable[str],
    font_size: int = 16,
    color: RGBColor = INK,
    font_name: str = BODY_FONT,
    bold_first: bool = False,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and index == 0
        run.font.color.rgb = color
        paragraph.space_after = Pt(5)
    return textbox


def add_section_title(slide, eyebrow: str, title: str, subtitle: str) -> None:
    add_text(slide, Inches(0.55), Inches(0.35), Inches(2.8), Inches(0.3), eyebrow, 14, ACCENT, True)
    add_text(slide, Inches(0.55), Inches(0.65), Inches(8.4), Inches(0.55), title, 26, INK, True, TITLE_FONT)
    add_text(slide, Inches(0.55), Inches(1.1), Inches(11.8), Inches(0.45), subtitle, 13, MUTED)


def add_footer(slide, text: str) -> None:
    add_text(slide, Inches(0.6), Inches(7.02), Inches(12.0), Inches(0.25), text, 10, MUTED)


def add_stat_chip(slide, left, top, width, label: str, value: str, color: RGBColor) -> None:
    add_rect(slide, left, top, width, Inches(0.72), color)
    add_text(slide, left + Inches(0.16), top + Inches(0.08), width - Inches(0.32), Inches(0.18), label, 10, PAPER, True)
    add_text(
        slide,
        left + Inches(0.16),
        top + Inches(0.27),
        width - Inches(0.32),
        Inches(0.3),
        value,
        18,
        PAPER,
        True,
        TITLE_FONT,
    )


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body_lines: Iterable[str],
    accent_color: RGBColor,
) -> None:
    add_outline_rect(slide, left, top, width, height, accent_color, PAPER)
    add_rect(slide, left, top, width, Inches(0.16), accent_color, radius=False)
    add_text(slide, left + Inches(0.16), top + Inches(0.22), width - Inches(0.32), Inches(0.3), title, 16, INK, True)
    add_paragraphs(
        slide,
        left + Inches(0.16),
        top + Inches(0.62),
        width - Inches(0.32),
        height - Inches(0.76),
        body_lines,
        font_size=13,
        color=MUTED,
    )


def add_picture_contain(slide, image_path: Path, left, top, width, height) -> None:
    if not image_path.exists():
        return
    with Image.open(image_path) as img:
        img_width, img_height = img.size
    box_ratio = width / height
    img_ratio = img_width / img_height
    if img_ratio > box_ratio:
        final_width = width
        final_height = width / img_ratio
        final_left = left
        final_top = top + (height - final_height) / 2
    else:
        final_height = height
        final_width = height * img_ratio
        final_top = top
        final_left = left + (width - final_width) / 2
    slide.shapes.add_picture(str(image_path), final_left, final_top, final_width, final_height)


def add_code_box(slide, left, top, width, height, title: str, lines: Iterable[str]) -> None:
    add_rect(slide, left, top, width, height, CODE_BG)
    add_text(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), Inches(0.25), title, 12, GOLD, True)
    add_paragraphs(
        slide,
        left + Inches(0.16),
        top + Inches(0.4),
        width - Inches(0.32),
        height - Inches(0.52),
        lines,
        font_size=11,
        color=PAPER,
        font_name=MONO_FONT,
    )


def add_connector(slide, x1, y1, x2, y2, color: RGBColor = ACCENT_DARK) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.8)


def build_cover(prs: Any, package_json: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)

    add_rect(slide, Inches(0), Inches(0), Inches(5.3), Inches(7.5), ACCENT_DARK, radius=False)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(4.1), Inches(0.35), "Shelf AI / VS Code Extension", 16, GOLD, True)
    add_text(
        slide,
        Inches(0.6),
        Inches(1.1),
        Inches(4.2),
        Inches(1.15),
        "当前插件的结构、能力\n以及它如何在后台约束 Codex 编程",
        28,
        PAPER,
        True,
        TITLE_FONT,
    )
    add_paragraphs(
        slide,
        Inches(0.6),
        Inches(2.55),
        Inches(4.0),
        Inches(1.25),
        [
            "不是启动器，不接管 Codex。",
            "它像结构版 ESLint，一直守在工作区背后。",
            "目标是把 Framework -> Product Spec -> Implementation Config -> Code -> Evidence 扣成一条线。",
        ],
        font_size=15,
        color=RGBColor(233, 233, 233),
    )

    stats = build_plugin_stats(package_json)

    add_stat_chip(slide, Inches(0.6), Inches(5.45), Inches(1.3), "版本", stats.version, ACCENT)
    add_stat_chip(slide, Inches(2.05), Inches(5.45), Inches(1.35), "命令", str(stats.command_count), ACCENT_ALT)
    add_stat_chip(slide, Inches(3.55), Inches(5.45), Inches(1.35), "配置", str(stats.config_count), GOLD)
    add_text(slide, Inches(0.6), Inches(6.45), Inches(4.2), Inches(0.25), "2026-03-10 · 针对当前仓库实现自动生成", 10, RGBColor(210, 214, 219))

    add_outline_rect(slide, Inches(5.65), Inches(0.6), Inches(7.1), Inches(3.75), RGBColor(222, 214, 199), PAPER)
    add_picture_contain(slide, ASSETS.hero_image, Inches(5.8), Inches(0.78), Inches(6.8), Inches(3.4))

    add_card(
        slide,
        Inches(5.75),
        Inches(4.65),
        Inches(2.15),
        Inches(1.55),
        "定位",
        ["无感安装", "后台检查", "主动命令可用"],
        ACCENT,
    )
    add_card(
        slide,
        Inches(8.15),
        Inches(4.65),
        Inches(2.15),
        Inches(1.55),
        "守卫对象",
        ["framework/*.md", "project spec/config", "generated/* 禁区"],
        ACCENT_ALT,
    )
    add_card(
        slide,
        Inches(10.55),
        Inches(4.65),
        Inches(2.15),
        Inches(1.55),
        "结果回流",
        ["Problems", "状态栏 / Sidebar", "pre-push hook"],
        GOLD,
    )


def build_positioning_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "01 / 定位",
        "Shelf AI 的核心定位：无感后台守卫，不是 Codex 启动器",
        "插件装上就开始检查；不装就不检查。它不接管 AI 进程，而是约束工作区里的真实变更。",
    )

    layers = [
        ("Framework", ACCENT, "规则与结构语言"),
        ("Product Spec", ACCENT_ALT, "产品真相"),
        ("Implementation Config", GOLD, "实现细化"),
        ("Code", ACCENT_DARK, "src/ scripts/ tests/"),
        ("Evidence", OK, "generated/* + diagnostics"),
    ]
    top = Inches(1.75)
    for index, (label, color, note) in enumerate(layers):
        y = top + Inches(0.82) * index
        add_rect(slide, Inches(0.85), y, Inches(3.6), Inches(0.6), color)
        add_text(slide, Inches(1.05), y + Inches(0.08), Inches(1.8), Inches(0.18), label, 18, PAPER, True, TITLE_FONT)
        add_text(slide, Inches(2.55), y + Inches(0.11), Inches(1.7), Inches(0.16), note, 11, PAPER)
        if index < len(layers) - 1:
            add_connector(slide, Inches(2.65), y + Inches(0.6), Inches(2.65), y + Inches(0.82))

    add_card(
        slide,
        Inches(5.0),
        Inches(1.8),
        Inches(3.2),
        Inches(1.7),
        "它做什么",
        [
            "监听保存、创建、重命名、删除、外部变更、窗口重新聚焦。",
            "把变化分类成：要不要物化、要不要跑 mypy、是否命中 generated 禁区。",
            "把结果送回 VS Code 诊断、状态栏、侧边栏与 git hook。",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(8.45),
        Inches(1.8),
        Inches(3.8),
        Inches(1.7),
        "它不做什么",
        [
            "不负责启动 Codex，也不劫持用户原来的使用入口。",
            "不靠聊天提示词“劝说” AI 守规矩。",
            "当前不是写前拦截，而是写后立即检查 + 推送前硬挡。",
        ],
        ACCENT_ALT,
    )

    add_code_box(
        slide,
        Inches(5.0),
        Inches(3.9),
        Inches(7.2),
        Inches(2.0),
        "现在的基本判断",
        [
            "装插件 -> 默默运行后台守卫",
            "改 framework/spec/config -> 自动 materialize",
            "改 src/scripts/tests 下 Python -> 追加 mypy",
            "改 projects/*/generated/* -> normal 报错 / strict 恢复",
            "推送前 -> .githooks 再做一次硬验证",
        ],
    )
    add_footer(slide, "结构约束来源：AGENTS.md；插件定位实现：tools/vscode/shelf-ai/extension.js 与 guarding.js")


def build_structure_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "02 / 结构",
        "插件由一个 orchestrator 加三个专职模块，再连到仓库脚本",
        "真正的“后台守卫链”不在单一命令里，而在 extension.js 对 VS Code 事件、providers、脚本和 hooks 的统一编排里。",
    )

    add_card(
        slide,
        Inches(0.7),
        Inches(1.8),
        Inches(2.5),
        Inches(1.2),
        "VS Code Surface",
        ["Activity Bar / Sidebar", "Status Bar / Problems", "Commands / File Events"],
        ACCENT,
    )
    add_card(
        slide,
        Inches(3.55),
        Inches(1.65),
        Inches(3.0),
        Inches(1.5),
        "extension.js",
        ["唯一 orchestrator", "注册 providers / commands / watchers", "执行 runValidation 与 UI 回流"],
        ACCENT_DARK,
    )
    add_card(
        slide,
        Inches(6.9),
        Inches(1.8),
        Inches(2.7),
        Inches(1.2),
        "guarding.js",
        ["变更分类器", "watched path / generated protect", "materialize + mypy 判定"],
        ACCENT_ALT,
    )
    add_card(
        slide,
        Inches(9.95),
        Inches(1.8),
        Inches(2.65),
        Inches(1.2),
        "仓库脚本链",
        ["materialize_project.py", "validate_strict_mapping.py", "uv run mypy"],
        GOLD,
    )

    add_connector(slide, Inches(3.2), Inches(2.4), Inches(3.55), Inches(2.4))
    add_connector(slide, Inches(6.55), Inches(2.4), Inches(6.9), Inches(2.4))
    add_connector(slide, Inches(9.6), Inches(2.4), Inches(9.95), Inches(2.4))

    add_card(
        slide,
        Inches(1.0),
        Inches(3.55),
        Inches(3.2),
        Inches(1.7),
        "framework_navigation.js",
        [
            "识别 framework markdown 符号。",
            "支持 Go to Definition / Hover / References。",
            "还能把 boundary token 追到 product_spec.toml section。",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.55),
        Inches(3.55),
        Inches(3.2),
        Inches(1.7),
        "framework_completion.js",
        [
            "管理 @framework 模板与骨架片段。",
            "按上下文返回 completion entries。",
            "保证框架作者入口始终直接可用。",
        ],
        ACCENT_ALT,
    )
    add_card(
        slide,
        Inches(8.1),
        Inches(3.55),
        Inches(3.2),
        Inches(1.7),
        "package.json / README / tests",
        [
            "贡献命令、视图、配置、激活事件。",
            "README 公开行为契约。",
            "test_snippets.js / test_guarding.js 防回归。",
        ],
        GOLD,
    )

    add_code_box(
        slide,
        Inches(0.9),
        Inches(5.7),
        Inches(11.8),
        Inches(0.95),
        "关键入口",
        [
            "registerWebviewViewProvider -> Sidebar；registerDefinitionProvider / Hover / Reference / Completion -> 框架作者体验；"
            "onDidSave / Create / Delete / Rename / Focus / FileSystemWatcher -> 后台守卫触发链",
        ],
    )
    add_footer(
        slide,
        "文件分工：extension.js:869-1279；framework_navigation.js；framework_completion.js；guarding.js；package.json",
    )


def build_capability_slide(prs: Any, package_json: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "03 / 能力",
        "当前插件已经覆盖“写、看、查、守、交付”五个层面",
        "对用户来说，最好用的地方不是多一个聊天框，而是把仓库结构语言直接变成 VS Code 的可见能力。",
    )

    add_card(
        slide,
        Inches(0.7),
        Inches(1.65),
        Inches(3.0),
        Inches(1.6),
        "写：框架作者入口",
        [
            "@framework 模板命令与 snippet 并存。",
            "章节、C/B/R/V、边界与规则子项自动补全。",
            "这是仓库硬契约，不能被重构掉。",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(0.7),
        Inches(3.45),
        Inches(3.0),
        Inches(1.6),
        "看：结构导航",
        [
            "侧边栏打开 framework tree。",
            "Ctrl/Cmd+Click、Hover、Find All References 全接上。",
            "boundary 还能反查 product_spec section。",
        ],
        ACCENT_ALT,
    )
    add_card(
        slide,
        Inches(0.7),
        Inches(5.25),
        Inches(3.0),
        Inches(1.55),
        "守：后台校验",
        [
            "启动自动校验，窗口回焦也会补跑。",
            "保存/外部变更后自动分类并执行守卫链。",
            "状态栏、Problems、Sidebar、Output 一起回流。",
        ],
        GOLD,
    )

    add_outline_rect(slide, Inches(4.1), Inches(1.75), Inches(8.55), Inches(4.95), RGBColor(223, 214, 201))
    add_picture_contain(slide, ASSETS.workbench_image, Inches(4.28), Inches(1.92), Inches(8.2), Inches(4.58))
    add_text(
        slide,
        Inches(4.28),
        Inches(6.48),
        Inches(8.0),
        Inches(0.18),
        "示意图：插件最终守住的是知识库工作区及其结构化产物，而不是另起一套“AI 专用编辑器”。",
        10,
        MUTED,
    )

    stats = build_plugin_stats(package_json)
    add_stat_chip(slide, Inches(9.2), Inches(0.65), Inches(1.0), "命令", str(stats.command_count), ACCENT)
    add_stat_chip(slide, Inches(10.35), Inches(0.65), Inches(1.0), "配置", str(stats.config_count), ACCENT_ALT)
    add_stat_chip(slide, Inches(11.5), Inches(0.65), Inches(1.0), "激活", "7", GOLD)
    add_footer(slide, "能力面：framework_navigation.js、framework_completion.js、extension.js、package.json")


def build_pipeline_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "04 / 工作流",
        "后台守卫链：先分类，再决定跑什么，而不是每次一股脑全跑",
        "这是性能和严格性之间的工程平衡点。只有相关变更才触发相关命令，避免把每次保存都变成一次重型 CI。",
    )

    steps = [
        ("1. 事件触发", "save / create / rename / delete / watcher / focus", ACCENT),
        ("2. classify", "watch 范围、generated 禁区、materialize / mypy 判定", ACCENT_ALT),
        ("3. generated 守卫", "normal 报告；strict 自动恢复", GOLD),
        ("4. auto materialize", "framework/spec/config 变化才跑", ACCENT_DARK),
        ("5. mypy", "仅 src/ scripts/ tests 下 Python 改动", OK),
        ("6. strict validate", "validate_strict_mapping.py / --check-changes", WARN),
        ("7. 结果回流", "Diagnostics / Status / Sidebar / Output / Hook", ACCENT),
    ]

    left = Inches(0.7)
    top = Inches(2.0)
    width = Inches(1.7)
    for index, (title, note, color) in enumerate(steps):
        x = left + Inches(1.8) * index
        add_rect(slide, x, top, width, Inches(1.08), color)
        add_text(slide, x + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.22), title, 13, PAPER, True)
        add_text(slide, x + Inches(0.12), top + Inches(0.35), width - Inches(0.24), Inches(0.56), note, 10, PAPER)
        if index < len(steps) - 1:
            add_connector(slide, x + width, top + Inches(0.54), x + Inches(1.8), top + Inches(0.54))

    add_card(
        slide,
        Inches(0.9),
        Inches(4.15),
        Inches(3.45),
        Inches(1.8),
        "strict mode 分支",
        [
            "命中 projects/*/generated/* 直改时，优先反推出 owning product_spec.toml。",
            "能恢复就立即 materialize 回去，并抑制刚生成目录的重复事件。",
            "恢复失败才把问题抛给用户。",
        ],
        ACCENT,
    )
    add_card(
        slide,
        Inches(4.65),
        Inches(4.15),
        Inches(3.45),
        Inches(1.8),
        "normal mode 分支",
        [
            "同样识别 generated 禁区，但不主动覆盖文件。",
            "问题会进 Diagnostics / Problems。",
            "适合更温和的团队推广默认值。",
        ],
        ACCENT_ALT,
    )
    add_card(
        slide,
        Inches(8.4),
        Inches(4.15),
        Inches(3.45),
        Inches(1.8),
        "噪音控制",
        [
            "OutputChannel 保存最近结果，不落地日志文件。",
            "generated 目录有短时事件抑制，避免物化后自触发风暴。",
            "只有相关变化才补跑对应命令。",
        ],
        GOLD,
    )
    add_footer(slide, "核心实现：extension.js:247-345, 1312-1448；guarding.js:118-186")


def build_modes_slide(prs: Any, package_json: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "05 / 配置",
        "现在的配置已经能把“提醒型守卫”和“强约束守卫”区分开",
        "配置层不是附属品。对团队推广来说，它决定插件是先提示、先教育，还是直接升级为强闸门。",
    )

    add_outline_rect(slide, Inches(0.75), Inches(1.7), Inches(5.8), Inches(4.8), RGBColor(222, 214, 201))
    add_rect(slide, Inches(0.75), Inches(1.7), Inches(2.9), Inches(0.55), ACCENT)
    add_rect(slide, Inches(3.65), Inches(1.7), Inches(2.9), Inches(0.55), ACCENT_ALT)
    add_text(slide, Inches(0.95), Inches(1.84), Inches(2.45), Inches(0.2), "Normal", 18, PAPER, True, TITLE_FONT)
    add_text(slide, Inches(3.85), Inches(1.84), Inches(2.45), Inches(0.2), "Strict", 18, PAPER, True, TITLE_FONT)

    rows = [
        ("generated 直改", "报告问题，不覆盖文件", "自动尝试恢复；恢复失败再报错"),
        ("auto materialize", "可开关，默认开", "可开关，默认开"),
        ("Python 改动", "相关文件才跑 mypy", "相关文件才跑 mypy"),
        ("git hooks", "提示安装", "提示安装 + pre-push 更关键"),
        ("适用场景", "先推广、先适应", "核心团队、严格仓库"),
    ]
    row_top = int(Inches(2.35))
    for title, normal_text, strict_text in rows:
        add_text(slide, Inches(0.95), row_top, Inches(1.5), Inches(0.2), title, 12, MUTED, True)
        add_text(slide, Inches(2.1), row_top, Inches(1.3), Inches(0.42), normal_text, 12, INK)
        add_text(slide, Inches(3.95), row_top, Inches(2.2), Inches(0.42), strict_text, 12, INK)
        row_top += int(Inches(0.78))

    stats = build_plugin_stats(package_json)
    add_code_box(
        slide,
        Inches(6.95),
        Inches(1.82),
        Inches(5.65),
        Inches(2.15),
        "关键配置",
        [
            "shelf.guardMode",
            "shelf.autoMaterialize",
            "shelf.runMypyOnPythonChanges",
            "shelf.protectGeneratedFiles",
            "shelf.promptInstallGitHooks",
            "shelf.materializeCommand",
            "shelf.typeCheckCommand",
        ],
    )
    add_card(
        slide,
        Inches(6.95),
        Inches(4.18),
        Inches(5.65),
        Inches(2.15),
        f"当前配置面（共 {stats.config_count} 项）",
        [
            "行为开关：save / auto fail / guard mode / auto materialize / mypy / generated protect / hooks prompt。",
            "命令覆盖：change validation、full validation、framework tree generate、materialize、type check。",
            "说明已经从“单命令插件”进化成“可配置守卫系统”。",
        ],
        GOLD,
    )
    add_footer(slide, "配置来源：tools/vscode/shelf-ai/package.json；行为说明：tools/vscode/shelf-ai/README.md")


def build_implementation_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "06 / 代码实现",
        "真正值得看的不是 UI，而是每个文件分别扛了什么职责",
        "这页可以直接拿来给团队做 code walkthrough，因为它基本对应当前代码结构。",
    )

    cards = [
        (
            "extension.js",
            ["运行中心。", "注册 providers / commands / watchers。", "runValidation 把分类、命令执行和 UI 回流串起来。"],
            ACCENT,
        ),
        (
            "guarding.js",
            ["路径归一化与 watched path 判定。", "classifyWorkspaceChanges 决定 materialize / mypy / generated protect。", "把“跑什么”从 orchestrator 里抽出去。"],
            ACCENT_ALT,
        ),
        (
            "framework_navigation.js",
            ["框架引用解析。", "模块与边界 token 跳转、hover、reference。", "boundary 可映射到 product_spec.toml section。"],
            GOLD,
        ),
        (
            "framework_completion.js",
            ["@framework 模板与补全定义。", "按上下文给出标题、章节、C/B/R/V、R*.1~R*.4。", "保障框架作者入口不丢失。"],
            OK,
        ),
        (
            "package.json / README / tests",
            ["命令、配置、激活事件、view 贡献。", "README 对外说明行为。", "test_snippets.js / test_guarding.js 守住回归。"],
            WARN,
        ),
    ]

    start_x = Inches(0.72)
    start_y = Inches(1.75)
    for index, (title, lines, color) in enumerate(cards):
        row = index // 2
        col = index % 2
        left = start_x + Inches(6.15) * col
        top = start_y + Inches(1.72) * row
        width = Inches(5.7)
        height = Inches(1.45)
        if index == 4:
            left = Inches(3.8)
            top = Inches(5.15)
            width = Inches(5.7)
        add_card(slide, left, top, width, height, title, lines, color)

    add_code_box(
        slide,
        Inches(8.35),
        Inches(1.85),
        Inches(4.0),
        Inches(2.25),
        "关键函数",
        [
            "runValidation(...)",
            "refreshGitHookStatus(...)",
            "buildMaterializeCommand(...)",
            "parseStageFailure(...)",
            "parseMypyResult(...)",
            "classifyWorkspaceChanges(...)",
        ],
    )
    add_footer(
        slide,
        "入口分布：extension.js:869-1279；guarding.js:1-202；framework_navigation.js；framework_completion.js；test_snippets.js；test_guarding.js",
    )


def build_codex_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "07 / Codex 场景",
        "它如何约束 Codex：不拦聊天入口，而是拦工作区里的真实结果",
        "这也是为什么插件适合规模化分发。团队不需要统一 AI 客户端，只需要统一工作区守卫。",
    )

    add_card(
        slide,
        Inches(0.72),
        Inches(1.72),
        Inches(4.9),
        Inches(1.5),
        "Codex / 人 都一样",
        [
            "无论是 AI 改文件还是人改文件，最后都会变成 VS Code 工作区变更事件。",
            "Shelf AI 对事件做同一套后台处理，不为某个编辑者单独开后门。",
        ],
        ACCENT,
    )

    flow_steps = [
        ("Codex 产出 patch 或直接改文件", ACCENT_DARK),
        ("VS Code 触发 save / watcher", ACCENT),
        ("Shelf AI 分类并补跑必要校验", ACCENT_ALT),
        ("问题回到 Problems / 状态栏", GOLD),
        ("推送前再被 .githooks 硬挡一次", WARN),
    ]
    base_y = Inches(3.55)
    for index, (text, color) in enumerate(flow_steps):
        left = Inches(0.95) + Inches(1.02) * index
        width = Inches(0.92) if index < len(flow_steps) - 1 else Inches(1.25)
        add_rect(slide, left, base_y, width, Inches(1.35), color)
        add_text(slide, left + Inches(0.08), base_y + Inches(0.12), width - Inches(0.16), Inches(1.0), text, 10, PAPER, True)
        if index < len(flow_steps) - 1:
            add_connector(slide, left + width, base_y + Inches(0.67), left + Inches(1.02), base_y + Inches(0.67))

    add_outline_rect(slide, Inches(6.2), Inches(1.72), Inches(6.2), Inches(4.95), RGBColor(222, 214, 201))
    add_picture_contain(slide, ASSETS.chat_image, Inches(6.35), Inches(1.9), Inches(5.9), Inches(3.8))
    add_text(
        slide,
        Inches(6.35),
        Inches(5.9),
        Inches(5.8),
        Inches(0.35),
        "示意图：它最终守住的是知识库行为、引用链和产物一致性，而不是给 Codex 再包一层花哨 UI。",
        10,
        MUTED,
    )
    add_footer(slide, "当前边界：写后立即检查 + 推送前硬挡；尚未接管 Codex 进程，也不是写前拦截。")


def build_boundary_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_section_title(
        slide,
        "08 / 边界与结论",
        "现在这版插件已经足够做仓库级后台守卫，但它的边界也应该说清楚",
        "这页适合给团队做预期管理：哪些事现在已经可靠，哪些事如果要更强，需要换成 patch gate 或 MCP 写入闸门。",
    )

    add_card(
        slide,
        Inches(0.82),
        Inches(1.8),
        Inches(3.7),
        Inches(3.8),
        "已经能稳定做到",
        [
            "框架作者入口、跳转、hover、references、completion 全接通。",
            "后台识别 framework/spec/config/Python/generated 变化并执行分级守卫。",
            "strict mode 下可自动恢复 generated 禁区。",
            "git hooks 缺失会提示，一键安装后把 pre-push 也扣上。",
            "不生成仓库日志垃圾，输出只留在 VS Code OutputChannel。",
        ],
        OK,
    )
    add_card(
        slide,
        Inches(4.82),
        Inches(1.8),
        Inches(3.7),
        Inches(3.8),
        "当前明确边界",
        [
            "不启动 Codex，不绑定某个 AI 客户端。",
            "不是写前拦截，第一笔文件字节仍可能先落盘。",
            "对 generated 的“强拦”只覆盖明确禁区；一般代码改动还是先检查后提示。",
            "真正的结构约束仍依赖仓库规范、物化链和 validate 脚本。",
        ],
        WARN,
    )
    add_card(
        slide,
        Inches(8.82),
        Inches(1.8),
        Inches(3.7),
        Inches(3.8),
        "如果要再强一步",
        [
            "把 Codex 写入改成 patch gate，再由 Shelf 审批后落盘。",
            "把 framework/spec/config 的引用链进一步编译成更细的强规则。",
            "把“工作区后台守卫”升级成“代码写入闸门”。",
            "但那已经不是今天这版插件的定位了。",
        ],
        ACCENT,
    )

    add_outline_rect(slide, Inches(0.82), Inches(6.0), Inches(11.7), Inches(0.72), RGBColor(222, 214, 201))
    add_text(
        slide,
        Inches(1.0),
        Inches(6.18),
        Inches(11.3),
        Inches(0.22),
        "一句话结论：Shelf AI 当前已经是“结构感知 + 后台守卫 + 推送前硬挡”的实用插件；如果未来要做到绝对强制，就该升级成写入闸门，而不是继续堆提示词。",
        14,
        INK,
        True,
    )
    add_picture_contain(slide, ASSETS.branding_image, Inches(10.45), Inches(0.35), Inches(2.2), Inches(1.1))
    add_footer(slide, "此 PPT 由 scripts/build_shelf_ai_plugin_presentation.py 基于当前仓库实现生成。")


def build_presentation() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    package_json = load_package_json()

    prs: Any = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_cover(prs, package_json)
    build_positioning_slide(prs)
    build_structure_slide(prs)
    build_capability_slide(prs, package_json)
    build_pipeline_slide(prs)
    build_modes_slide(prs, package_json)
    build_implementation_slide(prs)
    build_codex_slide(prs)
    build_boundary_slide(prs)

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output_path = build_presentation()
    print(f"[PASS] wrote {output_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
